from __future__ import annotations
from typing import List, Tuple, Optional, Iterable
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import pathlib

def chunk_text(text: str, max_chars: int = 4000) -> List[str]:
    chunks = []
    i = 0
    while i < len(text):
        chunks.append(text[i:i+max_chars])
        i += max_chars
    return chunks

def _normalize_pages(pages: Optional[Iterable[int]]) -> Optional[set[int]]:
    if pages is None:
        return None
    s = set(int(p) for p in pages if isinstance(p, (int,)) or (isinstance(p, str) and p.isdigit()))
    return s or None

def parse_pdf(path: str, pages: Optional[Iterable[int]] = None) -> List[Tuple[str, dict]]:
    p = pathlib.Path(path)
    reader = PdfReader(str(p))
    out = []
    wanted = _normalize_pages(pages)
    for i, page in enumerate(reader.pages, start=1):
        if wanted is not None and i not in wanted:
            continue
        t = page.extract_text() or ""
        for ch in chunk_text(t):
            out.append((ch, {"modality":"file","filetype":"pdf","page":i,"filename":p.name}))
    return out

def parse_docx(path: str) -> List[Tuple[str, dict]]:
    p = pathlib.Path(path)
    doc = DocxDocument(str(p))
    text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
    return [(ch, {"modality":"file","filetype":"docx","filename":p.name}) for ch in chunk_text(text)]

def parse_pptx(path: str) -> List[Tuple[str, dict]]:
    p = pathlib.Path(path)
    prs = Presentation(str(p))
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        joined = "\n".join(texts)
        for ch in chunk_text(joined):
            out.append((ch, {"modality":"file","filetype":"pptx","slide":i,"filename":p.name}))
    return out

def parse_text(path: str) -> List[Tuple[str, dict]]:
    p = pathlib.Path(path)
    text = p.read_text(encoding="utf-8", errors="ignore")
    return [(ch, {"modality":"file","filetype":"text","filename":p.name}) for ch in chunk_text(text)]
